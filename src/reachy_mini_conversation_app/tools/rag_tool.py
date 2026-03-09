"""RAG tool — indexes and searches PDF/PPTX lecture slides."""

import json
import numpy as np
from pathlib import Path
from typing import Any
import faiss
from openai import OpenAI

client = OpenAI()

# Global in-memory store (per session)
_index: Any = None
_chunks: list[str] = []
_source_name: str = ""

EMBED_MODEL = "text-embedding-3-small"
EMBED_DIM = 1536
TOP_K = 4


def _embed(texts: list[str]) -> np.ndarray:
    """Get embeddings from OpenAI."""
    resp = client.embeddings.create(model=EMBED_MODEL, input=texts)
    return np.array([r.embedding for r in resp.data], dtype="float32")


def ingest_document(file_path: str) -> str:
    """Parse PDF or PPTX and build FAISS index. Returns status string."""
    global _index, _chunks, _source_name

    path = Path(file_path)
    _source_name = path.name
    raw_chunks: list[str] = []

    if path.suffix.lower() == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                raw_chunks.append(f"[Page {i+1}]\n{text.strip()}")

    elif path.suffix.lower() in (".pptx", ".ppt"):
        from pptx import Presentation
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                raw_chunks.append(f"[Slide {i+1}]\n" + "\n".join(texts))
    else:
        return f"Unsupported file type: {path.suffix}. Please upload a PDF or PPTX file."

    if not raw_chunks:
        return "No text content found in document."

    # Build FAISS index
    _chunks = raw_chunks
    embeddings = _embed(raw_chunks)
    _index = faiss.IndexFlatL2(EMBED_DIM)
    _index.add(embeddings)

    return f"Document '{_source_name}' loaded: {len(_chunks)} sections indexed and ready."


def rag_search(query: str) -> str:
    """Search indexed document for relevant content."""
    if _index is None or not _chunks:
        return "No document uploaded yet. Please upload your lecture slides first."

    q_embed = _embed([query])
    distances, indices = _index.search(q_embed, TOP_K)

    results = []
    for idx in indices[0]:
        if idx < len(_chunks):
            results.append(_chunks[idx])

    if not results:
        return "No relevant content found in the uploaded document."

    context = "\n\n---\n\n".join(results)
    return f"Relevant content from '{_source_name}':\n\n{context}"


def get_status() -> dict:
    """Returns current RAG status."""
    return {
        "document_loaded": _index is not None,
        "source": _source_name,
        "chunks": len(_chunks),
    }


# Tool definition for OpenAI function calling
RAG_TOOL_DEFINITION = {
    "type": "function",
    "function": {
        "name": "rag_search",
        "description": (
            "Search the student's uploaded lecture slides or documents for relevant content. "
            "Use this when the student asks about specific topics from their course materials, "
            "or when you want to ground your explanation in their actual slide content."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": "What to search for in the lecture materials.",
                }
            },
            "required": ["query"],
        },
    },
}